import os
import uuid
import tempfile
from pathlib import Path
from typing import Annotated, TypedDict, cast
from supabase import create_client, Client

import fitz
from docx import Document as DocxDocument
from dotenv import load_dotenv
from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from google import genai
from google.genai import types
from pptx.api import Presentation
from pydantic import BaseModel
from qdrant_client import QdrantClient
from qdrant_client.conversions.common_types import Points
from qdrant_client.models import (
    Distance,
    FieldCondition,
    Filter,
    MatchValue,
    PayloadSchemaType,
    PointStruct,
    VectorParams,
)

_ = load_dotenv()


def _require_env(name: str) -> str:
    value = os.getenv(name)
    if not value:
        raise RuntimeError(f"Missing required environment variable: {name}")
    return value


# ── Config ────────────────────────────────────────────────────────────────────

GEMINI_API_KEY   = os.getenv("GEMINI_API_KEY")
QDRANT_URL       = os.getenv("QDRANT_URL")
QDRANT_API_KEY   = os.getenv("QDRANT_API_KEY")
SUPABASE_URL = _require_env("SUPABASE_URL")
SUPABASE_SERVICE_KEY = _require_env("SUPABASE_SERVICE_KEY")
supabase_client: Client = create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)
COLLECTION_NAME  = os.getenv("COLLECTION_NAME", "researchlens")
CHUNK_SIZE       = int(os.getenv("CHUNK_SIZE", "600"))
CHUNK_OVERLAP    = int(os.getenv("CHUNK_OVERLAP", "50"))
TOP_K            = int(os.getenv("TOP_K_RESULTS", "8"))
EMBED_MODEL      = "gemini-embedding-2"
EMBED_DIM        = 3072  # fixed output dim for text-embedding-004
CHAT_MODEL = "gemini-2.5-pro"

gemini_client = genai.Client(api_key=GEMINI_API_KEY)
qdrant = QdrantClient(url=QDRANT_URL, api_key=QDRANT_API_KEY)

# ── FastAPI setup ─────────────────────────────────────────────────────────────

app = FastAPI(title="ResearchLens API", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],   # tighten this to your Lovable domain in production
    allow_methods=["*"],
    allow_headers=["*"],
)

# ── Ensure Qdrant collection exists ───────────────────────────────────────────

def ensure_collection():
    existing = [c.name for c in qdrant.get_collections().collections]
    if COLLECTION_NAME not in existing:
        _ = qdrant.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE),
        )
    # Always ensure the index exists
    _ = qdrant.create_payload_index(
        COLLECTION_NAME, "document_id", PayloadSchemaType.KEYWORD
    )

ensure_collection()

# ── Helpers ───────────────────────────────────────────────────────────────────

class PageText(TypedDict):
    page: int
    text: str


class ChunkPayload(TypedDict):
    document_id: str
    filename: str
    page: int
    text: str


def _first_embedding_values(result: types.EmbedContentResponse) -> list[float]:
    if not result.embeddings or result.embeddings[0].values is None:
        raise HTTPException(status_code=502, detail="Embedding API returned no vector.")
    return result.embeddings[0].values


def _require_payload(payload: object | None) -> ChunkPayload:
    if not isinstance(payload, dict):
        raise HTTPException(status_code=500, detail="Search hit missing payload.")
    return cast(ChunkPayload, cast(object, payload))


def _supabase_rows(data: object | None) -> list[dict[str, object]]:
    if not isinstance(data, list):
        return []
    rows: list[dict[str, object]] = []
    for item in data:
        if isinstance(item, dict):
            rows.append(cast(dict[str, object], cast(object, item)))
    return rows


def _row_str(row: dict[str, object], key: str, default: str = "") -> str:
    value = row.get(key)
    return str(value) if value is not None else default


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping word-based chunks."""
    words = text.split()
    chunks: list[str] = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i : i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        i += chunk_size - overlap
    return chunks


def extract_pdf(path: str) -> list[PageText]:
    """Return list of {page, text} dicts from a PDF."""
    doc = fitz.open(path)
    pages: list[PageText] = []
    for i in range(doc.page_count):
        text = doc[i].get_text().strip()
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages


def extract_pptx(path: str) -> list[PageText]:
    """Return list of {page, text} dicts from a PPTX (page = slide number)."""
    prs = Presentation(path)
    slides: list[PageText] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append({"page": i, "text": " ".join(texts)})
    return slides

def extract_docx(path: str) -> list[PageText]:
    """Return list of {page, text} dicts from a DOCX file."""
    doc = DocxDocument(path)
    chunks: list[PageText] = []
    current: list[str] = []
    page = 1
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        current.append(text)
        # Treat every 10 paragraphs as a "page" for citation purposes
        if len(current) >= 10:
            chunks.append({"page": page, "text": " ".join(current)})
            page += 1
            current = []
    if current:
        chunks.append({"page": page, "text": " ".join(current)})
    return chunks


def embed_texts(texts: list[str]) -> list[list[float]]:
    embeddings = []
    for text in texts:
        result = gemini_client.models.embed_content(
            model=EMBED_MODEL,
            contents=text,
        )
        embeddings.append(_first_embedding_values(result))
    return embeddings


def embed_query(text: str) -> list[float]:
    result = gemini_client.models.embed_content(
        model=EMBED_MODEL,
        contents=text,
    )
    return _first_embedding_values(result)

# ── Pydantic models ────────────────────────────────────────────────────────────

class IngestResponse(BaseModel):
    document_id: str
    filename: str
    chunks_indexed: int


class Citation(BaseModel):
    source: str
    page: int | None = None
    snippet: str
    score: float


def _parse_citations(value: object) -> list[Citation]:
    if not isinstance(value, list):
        return []
    citations: list[Citation] = []
    for item in value:
        if not isinstance(item, dict):
            continue
        raw = cast(dict[str, object], cast(object, item))
        page_val = raw.get("page")
        score_val = raw.get("score")
        citations.append(
            Citation(
                source=_row_str(raw, "source"),
                page=page_val if isinstance(page_val, int) else None,
                snippet=_row_str(raw, "snippet"),
                score=float(score_val) if isinstance(score_val, (int, float)) else 0.0,
            )
        )
    return citations


class QueryRequest(BaseModel):
    document_id: str | None = None   # filter to one doc, or None = search all
    question: str


class QueryResponse(BaseModel):
    answer: str
    citations: list[Citation]
    confidence: str   # "high" | "medium" | "low"

class ChatHistoryItem(BaseModel):
    id: str
    document_id: str
    filename: str
    question: str
    answer: str
    citations: list[Citation]
    confidence: str
    created_at: str

class DocumentRecord(BaseModel):
    id: str
    user_id: str
    document_id: str
    filename: str
    chunks_indexed: int
    created_at: str


class UserDocumentsResponse(BaseModel):
    documents: list[DocumentRecord]

class ChatHistoryResponse(BaseModel):
    history: list[ChatHistoryItem]

# ── Routes ────────────────────────────────────────────────────────────────────

@app.get("/")
def health():
    return {"status": "ok", "service": "ResearchLens API"}


@app.post("/ingest", response_model=IngestResponse)
async def ingest(request: Request, file: Annotated[UploadFile, File()]):
    """
    Accept a PDF or PPTX file, extract text, chunk it,
    embed with Gemini, and store in Qdrant.
    """
    if not file.filename:
        raise HTTPException(status_code=400, detail="Filename is required.")

    filename = file.filename
    suffix = Path(filename).suffix.lower()
    if suffix not in {".pdf", ".pptx", ".ppt", ".docx"}:
        raise HTTPException(
            status_code=415,
            detail=f"Unsupported file type '{suffix}'. Upload PDF or PPTX.",
        )

    document_id = str(uuid.uuid4())

    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        _ = tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        if suffix == ".pdf":
            pages = extract_pdf(tmp_path)
        elif suffix == ".docx":
            pages = extract_docx(tmp_path)
        else:
            pages = extract_pptx(tmp_path)
    finally:
        os.unlink(tmp_path)

    if not pages:
        raise HTTPException(status_code=422, detail="No readable text found in file.")

    # Build chunks with metadata
    all_chunks: list[str] = []
    all_meta: list[ChunkPayload] = []
    for page_info in pages:
        for chunk in chunk_text(page_info["text"]):
            all_chunks.append(chunk)
            all_meta.append({
                "document_id": document_id,
                "filename": filename,
                "page": page_info["page"],
                "text": chunk,
            })

    if not all_chunks:
        raise HTTPException(status_code=422, detail="Could not extract any text chunks.")

    # Embed in batches of 100 (Gemini limit)
    batch_size = 10
    all_vectors = []
    for i in range(0, len(all_chunks), batch_size):
        batch = all_chunks[i : i + batch_size]
        vecs = embed_texts(batch)
        all_vectors.extend(vecs)

    # Upsert into Qdrant
    points = [
        PointStruct(
            id=str(uuid.uuid4()),
            vector=vec,
            payload=cast(dict[str, object], cast(object, meta)),
        )
        for vec, meta in zip(all_vectors, all_meta)
    ]
    _ = qdrant.upsert(collection_name=COLLECTION_NAME, points=cast(Points, points))

    # Save to Supabase if user_id provided
    user_id = request.headers.get("X-User-Id")
    if user_id:
        supabase_client.table("documents").insert({
            "user_id": user_id,
            "document_id": document_id,
            "filename": filename,
            "chunks_indexed": len(points),
        }).execute()

    return IngestResponse(
        document_id=document_id,
        filename=filename,
        chunks_indexed=len(points),
    )


@app.post("/query", response_model=QueryResponse)
async def query(request: Request, req: QueryRequest):
    """
    Embed the question, retrieve top-K chunks from Qdrant,
    and ask Gemini to answer with citations.
    """
    if not req.question.strip():
        raise HTTPException(status_code=400, detail="Question cannot be empty.")

    query_vec = embed_query(req.question)

    # Optional filter: restrict to one document
    search_filter = None
    if req.document_id:
        search_filter = Filter(
            must=[FieldCondition(key="document_id", match=MatchValue(value=req.document_id))]
        )

    hits = qdrant.search(
        collection_name=COLLECTION_NAME,
        query_vector=query_vec,
        limit=TOP_K,
        query_filter=search_filter,
        with_payload=True,
    )

    if not hits:
        return QueryResponse(
            answer="No relevant content found in the uploaded documents.",
            citations=[],
            confidence="low",
        )

    # Build context block for Gemini
    context_parts: list[str] = []
    for i, hit in enumerate(hits, start=1):
        p = _require_payload(hit.payload)
        loc = f"p.{p['page']}" if p.get("page") else "unknown page"
        context_parts.append(f"[{i}] {p['filename']} ({loc}):\n{p['text']}")
    context = "\n\n".join(context_parts)

    system_prompt = """You are ResearchLens, a precise research assistant.
Answer the user's question using ONLY the provided context excerpts.
Rules:
- Be concise and direct (2-4 sentences for simple questions, up to a short paragraph for complex ones).
- Cite sources inline using [1], [2] etc. matching the numbered excerpts.
- If the context does not contain enough information, say so clearly.
- Never fabricate facts or add information not present in the excerpts."""

    user_prompt = f"""Context excerpts:
{context}

Question: {req.question}

Answer (with inline citations):"""

    response = gemini_client.models.generate_content(
        model=CHAT_MODEL,
        contents=user_prompt,
        config=types.GenerateContentConfig(system_instruction=system_prompt),
    )
    answer_text = (response.text or "").strip()

    # Build citations list
    citations: list[Citation] = []
    for i, hit in enumerate(hits, start=1):
        p = _require_payload(hit.payload)
        # Only include citations that were actually referenced
        if f"[{i}]" in answer_text:
            text = str(p["text"])
            citations.append(Citation(
                source=p["filename"],
                page=p["page"],
                snippet=text[:200] + ("…" if len(text) > 200 else ""),
                score=round(hit.score, 3),
            ))

    # Confidence based on top hit score
    top_score = hits[0].score if hits else 0
    confidence = "high" if top_score > 0.82 else "medium" if top_score > 0.65 else "low"
    
    # Save to chat history if user_id provided
    user_id = request.headers.get("X-User-Id")
    if user_id:
        supabase_client.table("chat_history").insert({
            "user_id": user_id,
            "document_id": req.document_id or "all",
            "filename": _require_payload(hits[0].payload)["filename"] if hits else "unknown",
            "question": req.question,
            "answer": answer_text,
            "citations": [c.model_dump() for c in citations],
            "confidence": confidence,
        }).execute()

    return QueryResponse(
        answer=answer_text,
        citations=citations,
        confidence=confidence,
    )


@app.delete("/documents/{document_id}")
def delete_document(document_id: str):
    """Remove all vectors for a given document from Qdrant."""
    _ = qdrant.delete(
        collection_name=COLLECTION_NAME,
        points_selector=Filter(
            must=[FieldCondition(key="document_id", match=MatchValue(value=document_id))]
        ),
    )
    return {"deleted": document_id}


@app.get("/history/{user_id}", response_model=ChatHistoryResponse)
def get_chat_history(user_id: str, document_id: str | None = None):
    """Get chat history for a user, optionally filtered by document."""
    query = supabase_client.table("chat_history")\
        .select("*")\
        .eq("user_id", user_id)\
        .order("created_at", desc=True)\
        .limit(50)
    
    if document_id:
        query = query.eq("document_id", document_id)
    
    result = query.execute()
    history: list[ChatHistoryItem] = []
    for row in _supabase_rows(result.data):
        history.append(
            ChatHistoryItem(
                id=_row_str(row, "id"),
                document_id=_row_str(row, "document_id"),
                filename=_row_str(row, "filename"),
                question=_row_str(row, "question"),
                answer=_row_str(row, "answer"),
                citations=_parse_citations(row.get("citations")),
                confidence=_row_str(row, "confidence"),
                created_at=_row_str(row, "created_at"),
            )
        )
    return ChatHistoryResponse(history=history)


@app.get("/documents/{user_id}", response_model=UserDocumentsResponse)
def get_user_documents(user_id: str):
    """Get all documents uploaded by a user."""
    result = supabase_client.table("documents")\
        .select("*")\
        .eq("user_id", user_id)\
        .order("created_at", desc=True)\
        .execute()
    documents: list[DocumentRecord] = []
    for row in _supabase_rows(result.data):
        chunks_val = row.get("chunks_indexed")
        documents.append(
            DocumentRecord(
                id=_row_str(row, "id"),
                user_id=_row_str(row, "user_id"),
                document_id=_row_str(row, "document_id"),
                filename=_row_str(row, "filename"),
                chunks_indexed=int(chunks_val) if isinstance(chunks_val, int) else 0,
                created_at=_row_str(row, "created_at"),
            )
        )
    return UserDocumentsResponse(documents=documents)


@app.delete("/history/{user_id}/{history_id}")
def delete_history_item(user_id: str, history_id: str):
    """Delete a specific chat history item."""
    supabase_client.table("chat_history")\
        .delete()\
        .eq("id", history_id)\
        .eq("user_id", user_id)\
        .execute()
    return {"deleted": history_id} 
